from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt


def create_sample_image(path: Path) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    image = Image.new("RGB", (720, 420), color=(239, 245, 239))
    draw = ImageDraw.Draw(image)
    draw.rounded_rectangle((40, 40, 680, 380), radius=24, outline=(31, 111, 80), width=8)
    draw.text((70, 120), "PPT Extractor Demo", fill=(16, 56, 37))
    draw.text((70, 210), "Image block for parser / DOCX tests", fill=(60, 80, 68))
    image.save(path)
    return path


def build_sample_pptx(path: Path) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    image_path = create_sample_image(path.with_suffix(".png"))

    prs = Presentation()

    slide1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide1.shapes.title.text = "概率论导论"
    content = slide1.placeholders[1].text_frame
    content.clear()
    p0 = content.paragraphs[0]
    r0 = p0.add_run()
    r0.text = "样本空间（sample space）描述实验全部可能结果。"
    r0.font.size = Pt(20)
    p1 = content.add_paragraph()
    p1.level = 1
    r1 = p1.add_run()
    r1.text = "事件（event）是样本空间的子集。"
    r1.font.size = Pt(18)
    p2 = content.add_paragraph()
    p2.level = 1
    r2 = p2.add_run()
    r2.text = "概率函数需要满足非负性、规范性与可列可加性。"
    r2.font.size = Pt(18)

    slide1.shapes.add_picture(str(image_path), Inches(5.3), Inches(1.45), width=Inches(3.4))
    caption = slide1.shapes.add_textbox(Inches(5.3), Inches(4.95), Inches(3.4), Inches(0.5))
    caption_frame = caption.text_frame
    caption_frame.text = "图 1：概率树示意图"
    caption_frame.paragraphs[0].runs[0].font.size = Pt(12)

    footer = slide1.shapes.add_textbox(Inches(0.4), Inches(6.9), Inches(4.5), Inches(0.3))
    footer.text_frame.text = "Lecture 01 | 2026 Spring"
    footer.text_frame.paragraphs[0].runs[0].font.size = Pt(10)

    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(5), Inches(0.6))
    title_run = title_box.text_frame.paragraphs[0].add_run()
    title_run.text = "离散型随机变量"
    title_run.font.size = Pt(28)
    title_run.font.bold = True

    table_shape = slide2.shapes.add_table(3, 3, Inches(0.6), Inches(1.3), Inches(5.5), Inches(1.7))
    table = table_shape.table
    table.cell(0, 0).text = "变量"
    table.cell(0, 1).text = "取值"
    table.cell(0, 2).text = "概率"
    table.cell(1, 0).text = "X"
    table.cell(1, 1).text = "0"
    table.cell(1, 2).text = "0.4"
    table.cell(2, 0).text = "X"
    table.cell(2, 1).text = "1"
    table.cell(2, 2).text = "0.6"

    table_caption = slide2.shapes.add_textbox(Inches(0.8), Inches(3.15), Inches(5), Inches(0.4))
    table_caption.text_frame.text = "表 1：离散型随机变量分布示例"
    table_caption.text_frame.paragraphs[0].runs[0].font.size = Pt(12)

    shape = slide2.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(6.4),
        Inches(1.4),
        Inches(2.4),
        Inches(1.1),
    )
    shape.text_frame.text = "事件 A -> 事件 B"
    shape.text_frame.paragraphs[0].runs[0].font.size = Pt(18)

    prs.save(path)
    return path
