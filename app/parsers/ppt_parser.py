from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.builders.ir_builder import IRBuilder
from app.config import AppConfig
from app.models.ir_models import DocumentMetadata, SlideDocument, SlideIRDocument
from app.parsers.diagram_parser import build_diagram_block
from app.parsers.image_parser import build_image_block, export_picture
from app.parsers.table_parser import build_table_block
from app.parsers.text_parser import build_text_block

logger = logging.getLogger(__name__)

DIAGRAM_TEXT_TYPES = {
    getattr(MSO_SHAPE_TYPE, name)
    for name in ("AUTO_SHAPE", "CANVAS", "FREEFORM")
    if hasattr(MSO_SHAPE_TYPE, name)
}
DIAGRAM_FALLBACK_TYPES = {
    getattr(MSO_SHAPE_TYPE, name)
    for name in ("AUTO_SHAPE", "CANVAS", "FREEFORM", "LINE", "LINE_INVERSE", "BENT_UP_ARROW")
    if hasattr(MSO_SHAPE_TYPE, name)
}


class PPTParser:
    def __init__(self, config: AppConfig, asset_dir: Path):
        self.config = config
        self.asset_dir = asset_dir
        self.ir_builder = IRBuilder()

    def parse(self, input_path: Path) -> SlideIRDocument:
        prs = Presentation(str(input_path))
        core = prs.core_properties
        metadata = DocumentMetadata(
            source_path=str(input_path),
            source_name=input_path.name,
            slide_count=len(prs.slides),
            slide_width_pt=prs.slide_width.pt,
            slide_height_pt=prs.slide_height.pt,
            author=getattr(core, "author", None),
            company=getattr(core, "company", None),
        )
        slides: list[SlideDocument] = []
        for slide_index, slide in enumerate(prs.slides, start=1):
            slide_ir = SlideDocument(
                slide_index=slide_index,
                width_pt=prs.slide_width.pt,
                height_pt=prs.slide_height.pt,
                notes=self._extract_notes(slide),
            )
            slide_ir.blocks = self._parse_shapes(slide.shapes, slide_index)
            slides.append(self.ir_builder.refine_slide(slide_ir))
        return SlideIRDocument(metadata=metadata, slides=slides)

    def _extract_notes(self, slide: Any) -> str | None:
        if not self.config.export_notes:
            return None
        if not getattr(slide, "has_notes_slide", False):
            return None
        notes = slide.notes_slide.notes_text_frame.text.strip()
        return notes or None

    def _parse_shapes(self, shapes: Any, slide_index: int, parent_prefix: str = "") -> list[Any]:
        blocks = []
        for z_index, shape in enumerate(shapes, start=1):
            block_prefix = f"s{slide_index:03d}-{parent_prefix}{z_index:03d}"
            try:
                blocks.extend(self._parse_single_shape(shape, slide_index, z_index, block_prefix))
            except Exception as exc:  # noqa: BLE001
                logger.exception("Failed to parse shape %s on slide %s", shape.name, slide_index)
                blocks.append(
                    build_diagram_block(
                        shape=shape,
                        block_id=f"{block_prefix}-error",
                        z_index=z_index,
                        reason=f"shape parse failure: {exc}",
                        asset_stem=self.asset_dir
                        / f"slide_{slide_index:03d}"
                        / f"{block_prefix}-error",
                    )
                )
        return blocks

    def _parse_single_shape(
        self, shape: Any, slide_index: int, z_index: int, block_prefix: str
    ) -> list[Any]:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            return self._parse_shapes(shape.shapes, slide_index, parent_prefix=f"{block_prefix}-")
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_path = (
                self.asset_dir / f"slide_{slide_index:03d}" / f"{block_prefix}.{shape.image.ext}"
            )
            export_picture(shape, image_path)
            return [build_image_block(shape, f"{block_prefix}-image", z_index, image_path)]
        if getattr(shape, "has_table", False):
            return [build_table_block(shape, f"{block_prefix}-table", z_index)]
        if getattr(shape, "has_text_frame", False):
            text_block = build_text_block(shape, f"{block_prefix}-text", z_index)
            if text_block:
                if (
                    shape.shape_type in DIAGRAM_TEXT_TYPES
                    and not text_block.style.is_placeholder_title
                ):
                    return [
                        build_diagram_block(
                            shape=shape,
                            block_id=f"{block_prefix}-diagram",
                            z_index=z_index,
                            reason="shape downgraded to diagram-text block",
                            asset_stem=self.asset_dir
                            / f"slide_{slide_index:03d}"
                            / f"{block_prefix}-diagram",
                            text_override=text_block.content.get("text"),
                        )
                    ]
                return [text_block]
        if shape.shape_type in DIAGRAM_FALLBACK_TYPES:
            return [
                build_diagram_block(
                    shape=shape,
                    block_id=f"{block_prefix}-diagram",
                    z_index=z_index,
                    reason="unsupported vector object downgraded to diagram block",
                    asset_stem=self.asset_dir
                    / f"slide_{slide_index:03d}"
                    / f"{block_prefix}-diagram",
                )
            ]
        return []
